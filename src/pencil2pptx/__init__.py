#!/usr/bin/env python3
"""
pen2pptx — 将 Pencil .pen 文件转换为 PowerPoint .pptx

通过 Pencil MCP server 获取精确布局数据，渲染为原生 PPT 元素。
需要 Pencil 桌面应用正在运行。

用法:
    python pen2pptx.py input.pen [-o output.pptx] [--font-scale 0.73]
"""

from __future__ import annotations

import argparse
import asyncio
import json
import logging
import os
import subprocess
import sys
import tempfile
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from mcp import ClientSession, StdioServerParameters
from mcp.client.stdio import stdio_client
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

logger = logging.getLogger(__name__)

# ═══════════════════════════════════════════════════════════════
# 常量
# ═══════════════════════════════════════════════════════════════

PX_TO_EMU = 9525  # 1 像素 = 9525 EMU
DEFAULT_FONT_SCALE = 0.73  # Pencil px → PPT pt 的默认系数

# Pencil MCP server 路径（根据实际安装位置修改）
DEFAULT_PENCIL_CMD = str(
    Path(os.environ.get("LOCALAPPDATA", ""))
    / "Programs" / "Pencil" / "resources"
    / "app.asar.unpacked" / "out" / "mcp-server-windows-x64.exe"
)

# Pencil 桌面应用路径
DEFAULT_PENCIL_APP = str(
    Path(os.environ.get("LOCALAPPDATA", ""))
    / "Programs" / "Pencil" / "Pencil.exe"
)

PENCIL_STARTUP_TIMEOUT = 15  # 等待 Pencil 启动的最大秒数

# ═══════════════════════════════════════════════════════════════
# 数据模型
# ═══════════════════════════════════════════════════════════════


@dataclass
class LayoutNode:
    """合并了属性和布局信息的节点"""
    id: str = ""
    node_type: str = ""
    name: str = ""
    x: float = 0.0
    y: float = 0.0
    width: float = 0.0
    height: float = 0.0
    fill: str | None = None
    content: str = ""
    font_family: str | None = None
    font_size: float | None = None
    font_weight: str = "normal"
    font_style: str = "normal"
    text_align: str = "left"
    line_height: float = 1.2
    letter_spacing: float = 0.0
    corner_radius: float = 0.0
    opacity: float = 1.0
    stroke_color: str | None = None
    stroke_width: float = 0.0
    icon_font_name: str = ""
    icon_font_family: str = ""
    icon_image_path: str = ""
    has_image_fill: bool = False  # fill 为图片类型（{"type": "image", ...}）
    image_fill_path: str = ""  # 图片填充导出后的本地路径
    image_fill_url: str = ""  # 图片填充的原始 URL（用于本地路径 fallback）
    context: str = ""  # 节点的 context 属性，如 "image"/"img" 表示整体导出为图片
    context_image_path: str = ""  # context 为 image/img 时导出的图片路径
    path_image_path: str = ""  # path 类型节点导出的图片路径
    layout: str = ""  # frame 的布局方式：vertical / horizontal / none
    justify_content: str = ""  # 主轴对齐：center / flex-start / flex-end 等
    align_items: str = ""  # 交叉轴对齐：center / flex-start / flex-end 等
    padding_top: float = 0.0
    padding_right: float = 0.0
    padding_bottom: float = 0.0
    padding_left: float = 0.0
    gap: float = 0.0
    children: list[LayoutNode] = field(default_factory=list)
    raw_props: dict = field(default_factory=dict)


@dataclass
class PageData:
    """一个页面（幻灯片）的数据"""
    id: str = ""
    name: str = ""
    width: float = 960.0
    height: float = 540.0
    bg_image: str = ""  # 整页背景图路径（用于有图片填充的页面）
    nodes: list[LayoutNode] = field(default_factory=list)


# ═══════════════════════════════════════════════════════════════
# MCP 客户端 — 获取节点属性 + 计算后布局 + 导出图标
# ═══════════════════════════════════════════════════════════════

class PencilMcpClient:

    def __init__(self, pencil_cmd: str = DEFAULT_PENCIL_CMD,
                 pencil_app: str = DEFAULT_PENCIL_APP):
        self.pencil_cmd = pencil_cmd
        self.pencil_app = pencil_app
        self._image_fill_export_ok = True  # 首次失败后跳过后续尝试
        self._context_img_export_ok = True

    def _ensure_pencil_running(self) -> None:
        """检测 Pencil 是否运行，未运行则自动启动并等待就绪"""
        # 检查进程是否存在
        try:
            result = subprocess.run(
                ["tasklist", "/FI", "IMAGENAME eq Pencil.exe", "/NH"],
                capture_output=True, text=True, timeout=5,
            )
            if "Pencil.exe" in result.stdout:
                return  # 已在运行
        except Exception:
            pass

        # 启动 Pencil
        app_path = Path(self.pencil_app)
        if not app_path.exists():
            raise FileNotFoundError(f"Pencil 应用未找到: {app_path}")

        print("Pencil 未运行，正在启动...")
        subprocess.Popen(
            [str(app_path)],
            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
            creationflags=subprocess.DETACHED_PROCESS | subprocess.CREATE_NEW_PROCESS_GROUP
            if sys.platform == "win32" else 0,
        )

        # 等待 Pencil 进程出现
        for i in range(PENCIL_STARTUP_TIMEOUT):
            time.sleep(1)
            try:
                result = subprocess.run(
                    ["tasklist", "/FI", "IMAGENAME eq Pencil.exe", "/NH"],
                    capture_output=True, text=True, timeout=5,
                )
                if "Pencil.exe" in result.stdout:
                    # 额外等待让 Pencil 完全初始化（包括 MCP server）
                    print(f"Pencil 已启动，等待初始化...")
                    time.sleep(8)
                    return
            except Exception:
                pass
            print(f"  等待 Pencil 启动... ({i + 1}s)")

        raise TimeoutError(f"Pencil 启动超时（{PENCIL_STARTUP_TIMEOUT}s）")

    async def fetch_pages(self, pen_file: str, page_indices: list[int] | None = None) -> list[PageData]:
        """获取页面数据。page_indices 为从1开始的页码列表，None 表示全部。"""
        self._ensure_pencil_running()

        server_params = StdioServerParameters(
            command=self.pencil_cmd, args=["--app", "desktop"],
        )
        async with stdio_client(server_params) as (read, write):
            async with ClientSession(read, write) as session:
                await session.initialize()
                return await self._fetch_impl(session, pen_file, page_indices)

    async def _fetch_impl(self, session: ClientSession, pen_file: str, page_indices: list[int] | None = None) -> list[PageData]:
        # 统一使用绝对路径
        pen_file = str(Path(pen_file).resolve())

        # 确保文件在 Pencil 中打开（export_nodes 需要文件已打开）
        doc_opened = False
        for attempt in range(5):
            try:
                result = await session.call_tool("open_document", {"filePathOrTemplate": pen_file})
                for item in result.content:
                    if hasattr(item, "text"):
                        if "opened" in item.text.lower() or "document" in item.text.lower():
                            doc_opened = True
                if doc_opened:
                    print("  文件已在 Pencil 中打开")
                    await asyncio.sleep(2)
                    break
            except Exception as e:
                print(f"  打开文档重试 ({attempt + 1}/5): {e}")
                await asyncio.sleep(3)

        top_layout = await self._call(session, "snapshot_layout", {
            "filePath": pen_file, "maxDepth": 0,
        })

        if isinstance(top_layout, list):
            top_layout.sort(key=lambda f: (float(f.get("y", 0)), float(f.get("x", 0))))

        total = len(top_layout)
        print(f"  共 {total} 页")

        # 按页码筛选要处理的页面
        if page_indices:
            selected = [(i, top_layout[i - 1]) for i in page_indices if 1 <= i <= total]
            print(f"  选取第 {','.join(str(i) for i in page_indices)} 页")
        else:
            selected = [(i + 1, f) for i, f in enumerate(top_layout)]

        # 临时目录用于导出图片
        export_dir = tempfile.mkdtemp(prefix="pencil_export_")
        pages: list[PageData] = []
        process_total = len(selected)

        for seq, (page_num, frame) in enumerate(selected, 1):
            pid = frame["id"]

            # 1. 获取页面名称
            page_info = await self._call(session, "batch_get", {
                "filePath": pen_file, "nodeIds": [pid], "readDepth": 1,
            })
            page_name = page_info[0].get("name", "") if page_info else ""
            print(f"  [{seq}/{process_total}] 第{page_num}页 {page_name}")

            # 2. 获取布局
            print(f"    获取布局...")
            await asyncio.sleep(0.5)
            layout = await self._call(session, "snapshot_layout", {
                "filePath": pen_file, "parentId": pid, "maxDepth": 20,
            })

            # 3. 获取属性（带变量解析）
            print(f"    获取属性...")
            await asyncio.sleep(0.5)
            props_list = await self._call(session, "batch_get", {
                "filePath": pen_file, "parentId": pid,
                "readDepth": 10, "resolveInstances": True,
                "resolveVariables": True,
            })
            props_map = self._build_props_map(props_list)

            # 4. 合并节点
            nodes = [self._merge(c, props_map) for c in layout.get("children", []) if isinstance(c, dict)]
            node_count = _count(nodes)
            print(f"    {node_count} 个节点")

            # 5. 导出图片填充节点（has_image_fill）
            page_img_fills: list[LayoutNode] = []
            self._collect_image_fills(nodes, page_img_fills)
            if page_img_fills and doc_opened:
                print(f"    导出 {len(page_img_fills)} 个图片填充节点...")
                await asyncio.sleep(1)
                await self._export_image_fills(session, pen_file, page_img_fills)

            # 6. 导出 path 节点（SVG 路径无法用原生 PPT 还原）
            page_paths: list[LayoutNode] = []
            self._collect_paths(nodes, page_paths)
            if page_paths and doc_opened:
                print(f"    导出 {len(page_paths)} 个 path 节点...")
                await asyncio.sleep(1)
                await self._export_paths(session, pen_file, page_paths)

            # 7. 导出图标
            page_icons: list[LayoutNode] = []
            self._collect_icons(nodes, page_icons)
            if page_icons and doc_opened:
                print(f"    导出 {len(page_icons)} 个图标...")
                await asyncio.sleep(1)
                await self._export_icons(session, pen_file, page_icons)

            # 8. 导出 context 图片
            page_context_images: list[LayoutNode] = []
            self._collect_context_images(nodes, page_context_images)
            if page_context_images and doc_opened:
                print(f"    导出 {len(page_context_images)} 个 context 图片...")
                await asyncio.sleep(1)
                await self._export_context_images(session, pen_file, page_context_images)

            pages.append(PageData(
                id=pid, name=page_name,
                width=float(frame.get("width", 960)),
                height=float(frame.get("height", 540)),
                nodes=nodes,
            ))

            # 页间间隔，让 Pencil 有时间回收资源
            if seq < process_total:
                await asyncio.sleep(1)

        return pages

    # --- 图标导出 ---

    def _collect_icons(self, nodes: list[LayoutNode], out: list[LayoutNode]) -> None:
        for n in nodes:
            if n.node_type == "icon_font" and n.id:
                out.append(n)
            self._collect_icons(n.children, out)

    async def _export_icons(self, session: ClientSession, pen_file: str, icons: list[LayoutNode]) -> None:
        if not icons:
            return
        icon_dir = tempfile.mkdtemp(prefix="pencil_icons_")
        try:
            await asyncio.wait_for(
                session.call_tool("export_nodes", {
                    "filePath": pen_file, "nodeIds": [n.id for n in icons],
                    "outputDir": icon_dir, "format": "png", "scale": 4,
                }),
                timeout=30,
            )
            id_map = {n.id: n for n in icons}
            for fname in os.listdir(icon_dir):
                base = os.path.splitext(fname)[0]
                if base in id_map:
                    id_map[base].icon_image_path = os.path.join(icon_dir, fname)
        except (asyncio.TimeoutError, Exception) as e:
            logger.warning("图标导出跳过: %s", type(e).__name__)

    # --- path 节点收集与导出 ---

    def _collect_paths(self, nodes: list[LayoutNode], out: list[LayoutNode]) -> None:
        """递归收集 path 类型的节点（SVG 路径无法用原生 PPT 还原）"""
        for n in nodes:
            if n.node_type == "path" and n.id:
                out.append(n)
            self._collect_paths(n.children, out)

    async def _export_paths(self, session: ClientSession, pen_file: str, nodes: list[LayoutNode]) -> None:
        """将 path 类型节点导出为 PNG 图片"""
        if not nodes:
            return
        img_dir = tempfile.mkdtemp(prefix="pencil_paths_")
        try:
            await asyncio.wait_for(
                session.call_tool("export_nodes", {
                    "filePath": pen_file, "nodeIds": [n.id for n in nodes],
                    "outputDir": img_dir, "format": "png", "scale": 4,
                }),
                timeout=30,
            )
            id_map = {n.id: n for n in nodes}
            for fname in os.listdir(img_dir):
                base = os.path.splitext(fname)[0]
                if base in id_map:
                    id_map[base].path_image_path = os.path.join(img_dir, fname)
        except (asyncio.TimeoutError, Exception) as e:
            logger.warning("path 节点导出跳过: %s", type(e).__name__)

    # --- context 为 image/img 的节点收集与导出 ---

    def _collect_image_fills(self, nodes: list[LayoutNode], out: list[LayoutNode]) -> None:
        """递归收集 fill 为图片类型的节点"""
        for n in nodes:
            if n.has_image_fill and n.id:
                out.append(n)
            self._collect_image_fills(n.children, out)

    async def _export_image_fills(self, session: ClientSession, pen_file: str, nodes: list[LayoutNode]) -> None:
        """将 fill 为图片类型的节点逐个导出为 PNG 图片，失败时从本地路径 fallback"""
        if not nodes:
            return
        pen_dir = str(Path(pen_file).parent)
        img_dir = tempfile.mkdtemp(prefix="pencil_img_fill_")
        for n in nodes:
            try:
                await asyncio.wait_for(
                    session.call_tool("export_nodes", {
                        "filePath": pen_file, "nodeIds": [n.id],
                        "outputDir": img_dir, "format": "png", "scale": 2,
                    }),
                    timeout=30,
                )
                img_path = os.path.join(img_dir, f"{n.id}.png")
                if os.path.exists(img_path):
                    n.image_fill_path = img_path
                    continue
            except (asyncio.TimeoutError, Exception):
                pass
            # fallback: 从 image_fill_url 解析本地路径
            if n.image_fill_url:
                url = n.image_fill_url.lstrip("./")
                local_path = os.path.join(pen_dir, url)
                if os.path.exists(local_path):
                    n.image_fill_path = local_path
                    logger.info("图片填充 %s 使用本地路径: %s", n.id, local_path)
                else:
                    logger.warning("图片填充 %s 本地路径不存在: %s", n.id, local_path)

    def _collect_context_images(self, nodes: list[LayoutNode], out: list[LayoutNode]) -> None:
        """递归收集 context 为 image 或 img 的节点"""
        for n in nodes:
            if n.context.lower() in ("image", "img") and n.id:
                out.append(n)
            else:
                # 只有非 context-image 节点才需要递归子节点
                self._collect_context_images(n.children, out)

    async def _export_context_images(self, session: ClientSession, pen_file: str, nodes: list[LayoutNode]) -> None:
        """将 context 为 image/img 的节点导出为 PNG 图片"""
        if not nodes:
            return
        img_dir = tempfile.mkdtemp(prefix="pencil_ctx_img_")
        try:
            await asyncio.wait_for(
                session.call_tool("export_nodes", {
                    "filePath": pen_file, "nodeIds": [n.id for n in nodes],
                    "outputDir": img_dir, "format": "png", "scale": 2,
                }),
                timeout=30,
            )
            id_map = {n.id: n for n in nodes}
            for fname in os.listdir(img_dir):
                base = os.path.splitext(fname)[0]
                if base in id_map:
                    id_map[base].context_image_path = os.path.join(img_dir, fname)
        except (asyncio.TimeoutError, Exception) as e:
            self._context_img_export_ok = False
            logger.warning("context image 导出跳过: %s", type(e).__name__)

    # --- 数据构建 ---

    def _build_props_map(self, nodes) -> dict[str, dict]:
        result: dict[str, dict] = {}
        self._collect_props(nodes if isinstance(nodes, list) else [nodes], result)
        return result

    def _collect_props(self, nodes: list, result: dict[str, dict]) -> None:
        for node in nodes:
            if not isinstance(node, dict):
                continue
            nid = node.get("id", "")
            if nid:
                result[nid] = node
            children = node.get("children", [])
            if isinstance(children, list):
                self._collect_props(children, result)

    def _merge(self, layout: dict, props_map: dict[str, dict]) -> LayoutNode:
        nid = layout.get("id", "")
        p = props_map.get(nid, {})

        stroke_color, stroke_width = None, 0.0
        sr = p.get("stroke")
        if isinstance(sr, dict):
            fv = sr.get("fill")
            if isinstance(fv, str):
                stroke_color = fv
            th = sr.get("thickness", 0)
            if isinstance(th, (int, float)):
                stroke_width = float(th)
            elif isinstance(th, dict):
                # thickness 可能是 {"bottom": 1, "left": 1, ...}，取最大值
                vals = [float(v) for v in th.values() if isinstance(v, (int, float))]
                stroke_width = max(vals) if vals else 0.0
        elif isinstance(sr, str):
            stroke_color = sr

        # 解析 fill：可能是字符串颜色、dict 图片填充、或 dict 渐变
        raw_fill = p.get("fill")
        fill_color = None
        has_image_fill = False
        image_fill_url = ""
        if isinstance(raw_fill, str) and _is_valid_color(raw_fill):
            fill_color = raw_fill
        elif isinstance(raw_fill, dict):
            if raw_fill.get("type") == "image" and raw_fill.get("enabled", True):
                has_image_fill = True
                image_fill_url = raw_fill.get("url", "")
            else:
                # 渐变等复杂填充，尝试提取颜色
                stops = raw_fill.get("stops") or raw_fill.get("colors") or []
                if stops and isinstance(stops, list):
                    first = stops[0]
                    c = first.get("color", "") if isinstance(first, dict) else str(first)
                    if _is_valid_color(c):
                        fill_color = c

        # 解析 cornerRadius：可能是单个数值或数组
        raw_cr = p.get("cornerRadius", 0)
        if isinstance(raw_cr, list):
            corner_radius = max(float(x) for x in raw_cr) if raw_cr else 0.0
        else:
            corner_radius = _sf(raw_cr, 0.0)

        # 解析 padding：可能是单个数值或数组 [top, right, bottom, left] 或 [vertical, horizontal]
        raw_pad = p.get("padding", 0)
        pad_t = pad_r = pad_b = pad_l = 0.0
        if isinstance(raw_pad, (int, float)):
            pad_t = pad_r = pad_b = pad_l = float(raw_pad)
        elif isinstance(raw_pad, list):
            if len(raw_pad) == 1:
                pad_t = pad_r = pad_b = pad_l = float(raw_pad[0])
            elif len(raw_pad) == 2:
                pad_t = pad_b = float(raw_pad[0])
                pad_r = pad_l = float(raw_pad[1])
            elif len(raw_pad) >= 4:
                pad_t, pad_r, pad_b, pad_l = [float(x) for x in raw_pad[:4]]

        node = LayoutNode(
            id=nid, node_type=p.get("type", ""), name=p.get("name", ""),
            x=float(layout.get("x", 0)), y=float(layout.get("y", 0)),
            width=float(layout.get("width", 0)), height=float(layout.get("height", 0)),
            fill=fill_color, content=p.get("content", ""),
            font_family=p.get("fontFamily"), font_size=_sf(p.get("fontSize")),
            font_weight=str(p.get("fontWeight", "normal")),
            font_style=str(p.get("fontStyle", "normal")),
            text_align=p.get("textAlign", "left"),
            line_height=_sf(p.get("lineHeight"), 1.2),
            letter_spacing=_sf(p.get("letterSpacing"), 0.0),
            corner_radius=corner_radius,
            opacity=_sf(p.get("opacity"), 1.0),
            stroke_color=stroke_color, stroke_width=stroke_width,
            icon_font_name=p.get("iconFontName", ""),
            icon_font_family=p.get("iconFontFamily", ""),
            has_image_fill=has_image_fill,
            image_fill_url=image_fill_url,
            context=str(p.get("context", "")),
            layout=str(p.get("layout", "")),
            justify_content=str(p.get("justifyContent", "")),
            align_items=str(p.get("alignItems", "")),
            padding_top=pad_t, padding_right=pad_r,
            padding_bottom=pad_b, padding_left=pad_l,
            gap=_sf(p.get("gap"), 0.0),
            raw_props=p,
        )
        for cl in layout.get("children", []):
            if isinstance(cl, dict):
                node.children.append(self._merge(cl, props_map))
        return node

    async def _call(self, session: ClientSession, tool: str, args: dict) -> Any:
        result = await session.call_tool(tool, args)
        for item in result.content:
            if hasattr(item, "text"):
                return json.loads(item.text)
        return {}


# ═══════════════════════════════════════════════════════════════
# PPT 渲染器
# ═══════════════════════════════════════════════════════════════

def render_pages(pages: list[PageData], output: Path, font_scale: float) -> None:
    prs = Presentation()
    for page in pages:
        prs.slide_width = Emu(int(page.width * PX_TO_EMU))
        prs.slide_height = Emu(int(page.height * PX_TO_EMU))
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 逐节点渲染，按图形情况决定用原生元素还是图片
        for node in page.nodes:
            _render(slide, node, 0.0, 0.0, font_scale)
    prs.save(str(output))


def _render(slide, n: LayoutNode, px: float, py: float, fs: float) -> None:
    ax, ay = px + n.x, py + n.y

    # context 为 image/img 的节点直接作为图片插入，不递归渲染子节点
    if n.context.lower() in ("image", "img") and n.context_image_path and os.path.exists(n.context_image_path):
        l, t, w, h = _box(ax, ay, n.width, n.height)
        slide.shapes.add_picture(n.context_image_path, l, t, w, h)
        return

    # path 类型节点 — 用导出的图片
    if n.node_type == "path":
        if n.path_image_path and os.path.exists(n.path_image_path):
            l, t, w, h = _box(ax, ay, n.width, n.height)
            slide.shapes.add_picture(n.path_image_path, l, t, w, h)
        return

    # fill 为图片类型的节点 — 用导出的图片
    if n.has_image_fill:
        if n.image_fill_path and os.path.exists(n.image_fill_path):
            l, t, w, h = _box(ax, ay, n.width, n.height)
            pic = slide.shapes.add_picture(n.image_fill_path, l, t, w, h)
            if n.opacity < 1.0:
                _set_picture_opacity(pic, n.opacity)
        # 图片填充的子节点仍需渲染（如上面的文字）
        for c in n.children:
            _render(slide, c, ax, ay, fs)
        return

    t = n.node_type

    if t == "text":
        _text(slide, n, ax, ay, fs)
    elif t == "rectangle":
        _rect(slide, n, ax, ay)
    elif t == "ellipse":
        _ellipse(slide, n, ax, ay)
    elif t == "line":
        _line(slide, n, ax, ay)
    elif t in ("frame", "group"):
        _frame(slide, n, ax, ay, fs)
    elif t == "icon_font":
        _icon(slide, n, ax, ay)


def _text(slide, n: LayoutNode, ax: float, ay: float, fs: float) -> None:
    if not n.content:
        return
    left = Emu(int(ax * PX_TO_EMU))
    top = Emu(int(ay * PX_TO_EMU))
    w = Emu(int(n.width * PX_TO_EMU)) if n.width > 0 else Emu(200 * PX_TO_EMU)
    h = Emu(int(n.height * PX_TO_EMU)) if n.height > 0 else Emu(30 * PX_TO_EMU)

    tf = slide.shapes.add_textbox(left, top, w, h).text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)

    align_map = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}

    # 统一换行符：将 \v (\x0B) 和 \r\n 都转为 \n，再按 \n 分段
    text = n.content.replace("\v", "\n").replace("\r\n", "\n").replace("\r", "\n")

    for i, line in enumerate(text.split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align_map.get(n.text_align, PP_ALIGN.LEFT)
        run = para.add_run()
        run.text = line
        f = run.font
        if n.font_family:
            f.name = n.font_family
        if n.font_size and n.font_size > 0:
            f.size = Pt(n.font_size * fs)
        if n.fill and isinstance(n.fill, str) and _is_valid_color(n.fill):
            f.color.rgb = _rgb(n.fill)
        if n.font_weight in ("bold", "600", "700", "800", "900"):
            f.bold = True
        if n.font_style == "italic":
            f.italic = True
        if n.opacity < 1.0:
            _set_text_opacity(f, n.opacity)
        if n.line_height and n.line_height != 1.0:
            # python-pptx 的 line_spacing 接受浮点数作为行距倍数
            para.line_spacing = n.line_height


def _rect(slide, n: LayoutNode, ax: float, ay: float) -> None:
    l, t, w, h = _box(ax, ay, n.width, n.height)
    shape_type = _shape_type_for(n)
    s = slide.shapes.add_shape(shape_type, l, t, w, h)
    if shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        _set_corner_radius(s, n.corner_radius)
    _fill(s, n); _stroke(s, n); _no_shadow(s)
    _set_opacity(s, n.opacity)


def _ellipse(slide, n: LayoutNode, ax: float, ay: float) -> None:
    l, t, w, h = _box(ax, ay, n.width, n.height)
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    _fill(s, n); _stroke(s, n); _no_shadow(s)
    _set_opacity(s, n.opacity)


def _line(slide, n: LayoutNode, ax: float, ay: float) -> None:
    l, t, w, h = _box(ax, ay, n.width, max(n.height, 0.5))
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    _fill(s, n); s.line.fill.background(); _no_shadow(s)
    _set_opacity(s, n.opacity)


def _shape_type_for(n: LayoutNode) -> int:
    """根据 cornerRadius 和尺寸选择合适的形状类型"""
    if n.corner_radius <= 0:
        return MSO_SHAPE.RECTANGLE
    short_side = min(n.width, n.height)
    # cornerRadius >= 短边一半时，视为圆形/椭圆
    if short_side > 0 and n.corner_radius >= short_side / 2:
        return MSO_SHAPE.OVAL
    return MSO_SHAPE.ROUNDED_RECTANGLE


def _frame(slide, n: LayoutNode, ax: float, ay: float, fs: float) -> None:
    # 绘制 frame 自身的背景/边框
    if n.fill or (n.stroke_color and _is_valid_color(n.stroke_color) and n.stroke_width > 0):
        l, t, w, h = _box(ax, ay, n.width, n.height)
        shape_type = _shape_type_for(n)
        s = slide.shapes.add_shape(shape_type, l, t, w, h)
        if shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
            _set_corner_radius(s, n.corner_radius)
        _fill(s, n); _stroke(s, n); _no_shadow(s)
        _set_opacity(s, n.opacity)

    # 检测：frame 有居中对齐且唯一子节点是文本 → 用 PPT 原生居中
    text_children = [c for c in n.children if c.node_type == "text" and c.content]
    non_text_children = [c for c in n.children if c.node_type != "text" or not c.content]
    h_center = n.align_items == "center" or n.justify_content == "center"

    if h_center and len(text_children) == 1 and len(non_text_children) == 0:
        _text_in_frame(slide, text_children[0], n, ax, ay, fs)
        return

    for c in n.children:
        _render(slide, c, ax, ay, fs)


def _text_in_frame(slide, tn: LayoutNode, frame: LayoutNode, ax: float, ay: float, fs: float) -> None:
    """在居中对齐的 frame 内渲染文本，使用 PPT 原生对齐代替精确坐标"""
    from pptx.enum.text import MSO_ANCHOR

    # 文本框覆盖整个 frame 区域
    left = Emu(int(ax * PX_TO_EMU))
    top = Emu(int(ay * PX_TO_EMU))
    w = Emu(int(frame.width * PX_TO_EMU))
    h = Emu(int(frame.height * PX_TO_EMU))

    tf = slide.shapes.add_textbox(left, top, w, h).text_frame
    tf.word_wrap = True
    tf.auto_size = None

    # padding
    tf.margin_left = Emu(int(frame.padding_left * PX_TO_EMU))
    tf.margin_right = Emu(int(frame.padding_right * PX_TO_EMU))
    tf.margin_top = Emu(int(frame.padding_top * PX_TO_EMU))
    tf.margin_bottom = Emu(int(frame.padding_bottom * PX_TO_EMU))

    # 垂直对齐：justifyContent 控制 vertical layout 的主轴
    jc = frame.justify_content
    ai = frame.align_items
    if frame.layout == "vertical":
        # vertical layout: justifyContent 控制垂直方向
        if jc == "center":
            tf.paragraphs[0].space_before = tf.paragraphs[0].space_after = None
            tf.word_wrap = True
            # 设置垂直居中
            txBody = tf._txBody
            from lxml import etree
            ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
            bodyPr = txBody.find(f"{{{ns}}}bodyPr")
            if bodyPr is not None:
                bodyPr.set("anchor", "ctr")
        # alignItems 控制水平方向
        if ai == "center":
            h_align = PP_ALIGN.CENTER
        else:
            h_align = PP_ALIGN.LEFT
    else:
        # horizontal layout: justifyContent 控制水平方向
        if jc == "center":
            h_align = PP_ALIGN.CENTER
        else:
            h_align = PP_ALIGN.LEFT
        if ai == "center":
            txBody = tf._txBody
            from lxml import etree
            ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
            bodyPr = txBody.find(f"{{{ns}}}bodyPr")
            if bodyPr is not None:
                bodyPr.set("anchor", "ctr")
        h_align = h_align if jc == "center" or ai == "center" else PP_ALIGN.LEFT

    # 渲染文本内容
    text = tn.content.replace("\v", "\n").replace("\r\n", "\n").replace("\r", "\n")
    for i, line in enumerate(text.split("\n")):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = h_align
        run = para.add_run()
        run.text = line
        f = run.font
        if tn.font_family:
            f.name = tn.font_family
        if tn.font_size and tn.font_size > 0:
            f.size = Pt(tn.font_size * fs)
        if tn.fill and isinstance(tn.fill, str) and _is_valid_color(tn.fill):
            f.color.rgb = _rgb(tn.fill)
        if tn.font_weight in ("bold", "600", "700", "800", "900"):
            f.bold = True
        if tn.font_style == "italic":
            f.italic = True
        if tn.opacity < 1.0:
            _set_text_opacity(f, tn.opacity)
        if tn.line_height and tn.line_height != 1.0:
            para.line_spacing = tn.line_height


def _icon(slide, n: LayoutNode, ax: float, ay: float) -> None:
    if not n.icon_image_path or not os.path.exists(n.icon_image_path):
        return
    l, t, w, h = _box(ax, ay, n.width, n.height)
    slide.shapes.add_picture(n.icon_image_path, l, t, w, h)


# --- 工具函数 ---

def _no_shadow(shape):
    """移除形状的默认阴影效果"""
    from lxml import etree
    sp = shape._element
    nsmap = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    }
    # 方法1：将 style 中的 effectRef idx 设为 0（禁用主题效果引用）
    for eref in sp.findall(".//a:effectRef", nsmap):
        eref.set("idx", "0")
    # 方法2：在 spPr 中添加空 effectLst 显式覆盖
    spPr = sp.find(".//p:spPr", nsmap)
    if spPr is None:
        spPr = sp.find(".//{%s}spPr" % nsmap["a"])
    if spPr is not None:
        ns_a = nsmap["a"]
        for el in spPr.findall("{%s}effectLst" % ns_a):
            spPr.remove(el)
        etree.SubElement(spPr, "{%s}effectLst" % ns_a)


def _set_opacity(shape, opacity: float):
    """设置形状的整体透明度（0.0 完全透明 ~ 1.0 完全不透明）"""
    if opacity >= 1.0:
        return
    from lxml import etree
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    alpha_val = int(opacity * 100000)
    # 设置填充的 alpha
    for sf in shape._element.iter("{%s}solidFill" % ns_a):
        srgb = sf.find("{%s}srgbClr" % ns_a)
        if srgb is not None:
            # 移除已有的 alpha
            for a in srgb.findall("{%s}alpha" % ns_a):
                srgb.remove(a)
            alpha_el = etree.SubElement(srgb, "{%s}alpha" % ns_a)
            alpha_el.set("val", str(alpha_val))
    # 设置线条的 alpha
    for ln in shape._element.iter("{%s}ln" % ns_a):
        for sf in ln.iter("{%s}solidFill" % ns_a):
            srgb = sf.find("{%s}srgbClr" % ns_a)
            if srgb is not None:
                for a in srgb.findall("{%s}alpha" % ns_a):
                    srgb.remove(a)
                alpha_el = etree.SubElement(srgb, "{%s}alpha" % ns_a)
                alpha_el.set("val", str(alpha_val))


def _set_picture_opacity(pic_shape, opacity: float):
    """设置图片的透明度（通过 alphaModFix 实现）"""
    if opacity >= 1.0:
        return
    from lxml import etree
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    alpha_val = int(opacity * 100000)
    # 找到 blipFill 中的 blip 元素
    ns_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    for blip in pic_shape._element.iter("{%s}blip" % ns_a):
        # 移除已有的 alphaModFix
        for existing in blip.findall("{%s}alphaModFix" % ns_a):
            blip.remove(existing)
        alpha_el = etree.SubElement(blip, "{%s}alphaModFix" % ns_a)
        alpha_el.set("amt", str(alpha_val))


def _set_text_opacity(run_font, opacity: float):
    """设置文本 run 的透明度"""
    if opacity >= 1.0:
        return
    from lxml import etree
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    alpha_val = int(opacity * 100000)
    # 在 run 的 solidFill 中设置 alpha
    rPr = run_font._element if hasattr(run_font, '_element') else None
    if rPr is None:
        return
    for sf in rPr.iter("{%s}solidFill" % ns_a):
        srgb = sf.find("{%s}srgbClr" % ns_a)
        if srgb is not None:
            for a in srgb.findall("{%s}alpha" % ns_a):
                srgb.remove(a)
            alpha_el = etree.SubElement(srgb, "{%s}alpha" % ns_a)
            alpha_el.set("val", str(alpha_val))


def _set_corner_radius(shape, radius_px: float):
    """设置形状的圆角半径（Pencil px → EMU）"""
    from lxml import etree
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    spPr = shape._element.find(".//{%s}prstGeom" % ns_a)
    if spPr is not None:
        avLst = spPr.find("{%s}avLst" % ns_a)
        if avLst is None:
            avLst = etree.SubElement(spPr, "{%s}avLst" % ns_a)
        # 清除已有的 adj
        for gd in list(avLst):
            avLst.remove(gd)
        # PPT 圆角值是 1/50000 比例，基于短边
        short_side = min(shape.width, shape.height)
        if short_side > 0:
            radius_emu = int(radius_px * PX_TO_EMU)
            # adj 值 = (radius / 短边) * 50000，最大 50000
            adj_val = min(int(radius_emu / short_side * 50000), 50000)
            gd = etree.SubElement(avLst, "{%s}gd" % ns_a)
            gd.set("name", "adj")
            gd.set("fmla", f"val {adj_val}")

def _box(ax, ay, w, h):
    return Emu(int(ax * PX_TO_EMU)), Emu(int(ay * PX_TO_EMU)), Emu(int(w * PX_TO_EMU)), Emu(int(h * PX_TO_EMU))

def _fill(s, n: LayoutNode):
    if n.fill and isinstance(n.fill, str) and _is_valid_color(n.fill):
        s.fill.solid(); s.fill.fore_color.rgb = _rgb(n.fill)
    else:
        s.fill.background()

def _stroke(s, n: LayoutNode):
    if n.stroke_color and _is_valid_color(n.stroke_color) and n.stroke_width > 0:
        s.line.color.rgb = _rgb(n.stroke_color)
        # Pencil stroke thickness 是像素，转为 EMU（1px = 9525 EMU）
        s.line.width = Emu(int(n.stroke_width * PX_TO_EMU))
    else:
        s.line.fill.background()

def _is_valid_color(c) -> bool:
    """检查是否为有效的、非透明的十六进制颜色字符串"""
    if not isinstance(c, str):
        return False
    if c.startswith("$"):
        return False  # 变量引用，非颜色值
    h = c.lstrip("#")
    if len(h) not in (3, 6, 8):
        return False
    try:
        int(h, 16)
    except ValueError:
        return False
    # 8 位 RGBA：检查 alpha 通道，00 表示完全透明
    if len(h) == 8:
        alpha = int(h[6:8], 16)
        if alpha == 0:
            return False
    return True


def _rgb(hex_color: str) -> RGBColor:
    if not isinstance(hex_color, str):
        return RGBColor(0, 0, 0)
    h = hex_color.lstrip("#")
    if len(h) == 6:
        try:
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
        except ValueError:
            return RGBColor(0, 0, 0)
    elif len(h) == 8:
        # RGBA 格式，忽略 alpha
        try:
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
        except ValueError:
            return RGBColor(0, 0, 0)
    elif len(h) == 3:
        try:
            return RGBColor(int(h[0]*2, 16), int(h[1]*2, 16), int(h[2]*2, 16))
        except ValueError:
            return RGBColor(0, 0, 0)
    return RGBColor(0, 0, 0)

def _sf(v, d: float = 0.0) -> float:
    if v is None: return d
    try: return float(v)
    except (ValueError, TypeError): return d


# ═══════════════════════════════════════════════════════════════
# CLI 入口
# ═══════════════════════════════════════════════════════════════

def _parse_pages(s: str) -> list[int]:
    """解析页码字符串，如 '1,2,5' 或 '1-3,5'（页码从1开始）"""
    result = []
    for part in s.split(","):
        part = part.strip()
        if "-" in part:
            a, b = part.split("-", 1)
            result.extend(range(int(a), int(b) + 1))
        else:
            result.append(int(part))
    return sorted(set(result))


def main():
    ap = argparse.ArgumentParser(
        prog="pencil2pptx",
        description="将 Pencil .pen 文件转换为 PowerPoint .pptx（自动启动 Pencil 应用）",
    )
    ap.add_argument("input", help="输入 .pen 文件路径")
    ap.add_argument("-o", "--output", default=None, help="输出 .pptx 路径（默认同名 .pptx）")
    ap.add_argument("--pages", default=None,
                     help="导出指定页码，如 1,2 或 1-3,5（从1开始，默认全部）")
    ap.add_argument("--font-scale", type=float, default=DEFAULT_FONT_SCALE,
                     help=f"字体缩放系数，Pencil px → PPT pt（默认 {DEFAULT_FONT_SCALE}）")
    ap.add_argument("--pencil-cmd", default=DEFAULT_PENCIL_CMD,
                     help="Pencil MCP server 可执行文件路径")
    ap.add_argument("--pencil-app", default=DEFAULT_PENCIL_APP,
                     help="Pencil 桌面应用路径（用于自动启动）")
    parsed = ap.parse_args()

    inp = Path(parsed.input).resolve()
    if not inp.exists():
        print(f"错误: 文件不存在: {inp}", file=sys.stderr); sys.exit(1)

    out = Path(parsed.output) if parsed.output else inp.with_suffix(".pptx")

    try:
        client = PencilMcpClient(parsed.pencil_cmd, parsed.pencil_app)
        print(f"连接 Pencil MCP server...")

        # 解析页码
        page_indices = _parse_pages(parsed.pages) if parsed.pages else None
        pages = asyncio.run(client.fetch_pages(str(inp), page_indices))
        print(f"获取到 {len(pages)} 页")

        total_nodes = sum(_count(p.nodes) for p in pages)
        print(f"共 {total_nodes} 个节点")

        render_pages(pages, out, parsed.font_scale)
        print(f"完成: {out}")
    except Exception as e:
        print(f"失败: {e}", file=sys.stderr); sys.exit(1)


def _count(nodes) -> int:
    return sum(1 + _count(n.children) for n in nodes)


if __name__ == "__main__":
    main()
